import os
import re
import csv
from datetime import datetime, date, timedelta
from sqlalchemy import create_engine
from sqlalchemy.ext.declarative import declarative_base
from sqlalchemy.orm import sessionmaker, scoped_session
from config import DB_PATH

engine = create_engine(f'sqlite:///{DB_PATH}', connect_args={"check_same_thread": False})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Session = scoped_session(sessionmaker(bind=engine))
Base = declarative_base()

def add_event(subject, topic, activity, duration, note=''):
    from models import StudyEvent, KnowledgeNode, KnowledgePoint
    from config import REVIEW_INTERVALS
    session = Session()
    event = StudyEvent(date=datetime.now().date(), subject=subject, topic=topic, activity=activity, duration=duration, note=note)
    session.add(event)
    node = session.query(KnowledgeNode).filter_by(subject=subject, topic=topic).first()
    if not node:
        node = KnowledgeNode(subject=subject, topic=topic, last_review=datetime.now().date())
        session.add(node)
    node.last_review = datetime.now().date()
    today = date.today()
    for interval in REVIEW_INTERVALS:
        candidate = today + timedelta(days=interval)
        if candidate > today:
            node.next_review = candidate
            break
    else:
        node.next_review = today + timedelta(days=REVIEW_INTERVALS[-1])
    kp = session.query(KnowledgePoint).filter_by(subject=subject, status='未学').first()
    if kp and (topic in kp.title or kp.title in topic):
        kp.status = '已学'
    session.commit()
    session.close()
    from trainer import train_model

def get_events(days=90):
    from models import StudyEvent
    session = Session()
    since = date.today() - timedelta(days=days)
    res = session.query(StudyEvent).filter(StudyEvent.date >= since).all()
    session.close()
    return res

def get_all_events():
    from models import StudyEvent
    session = Session()
    res = session.query(StudyEvent).all()
    session.close()
    return res

def get_nodes(subject=None):
    from models import KnowledgeNode
    session = Session()
    q = session.query(KnowledgeNode)
    if subject:
        q = q.filter_by(subject)
    res = q.all()
    session.close()
    return res

def set_node_mastery(topic, mastery):
    from models import KnowledgeNode
    session = Session()
    node = session.query(KnowledgeNode).filter_by(topic=topic).first()
    if node:
        node.mastery = mastery
        session.commit()
    session.close()

def get_goals():
    from models import UserGoal
    session = Session()
    res = session.query(UserGoal).all()
    session.close()
    return res

def set_goal(subject, one_round_end, practice_start, current_phase='', phase_start=None, phase_end=None):
    from models import UserGoal
    session = Session()
    g = session.query(UserGoal).filter_by(subject=subject).first()
    if not g:
        g = UserGoal(subject=subject)
        session.add(g)
    g.one_round_end = datetime.strptime(one_round_end, '%Y-%m-%d').date()
    g.practice_start = datetime.strptime(practice_start, '%Y-%m-%d').date()
    if current_phase:
        g.current_phase = current_phase
    if phase_start:
        g.phase_start = datetime.strptime(phase_start, '%Y-%m-%d').date()
    if phase_end:
        g.phase_end = datetime.strptime(phase_end, '%Y-%m-%d').date()
    session.commit()
    session.close()

def save_daily_plan(date_str, plan_text):
    from models import DailyPlan
    session = Session()
    d = datetime.strptime(date_str, '%Y-%m-%d').date()
    p = session.query(DailyPlan).filter_by(date=d).first()
    if not p:
        p = DailyPlan(date=d)
        session.add(p)
    p.plan_text = plan_text
    session.commit()
    session.close()

def mark_plan_complete(date_str, completed, rate=None):
    from models import DailyPlan
    session = Session()
    d = datetime.strptime(date_str, '%Y-%m-%d').date()
    p = session.query(DailyPlan).filter_by(date=d).first()
    if not p:
        p = DailyPlan(date=d)
        session.add(p)
    p.completed = completed
    if rate is not None:
        p.completion_rate = rate
    session.commit()
    session.close()

def get_plan(date_str):
    from models import DailyPlan
    session = Session()
    d = datetime.strptime(date_str, '%Y-%m-%d').date()
    p = session.query(DailyPlan).filter_by(date=d).first()
    session.close()
    return p

def save_sleep(bed_time_str, wake_time_str, nap_duration, efficiency=0.0):
    from models import SleepRecord
    session = Session()
    d = date.today()
    rec = session.query(SleepRecord).filter_by(date=d).first()
    if not rec:
        rec = SleepRecord(date=d)
        session.add(rec)
    rec.bed_time = datetime.strptime(bed_time_str, '%H:%M').time()
    rec.wake_time = datetime.strptime(wake_time_str, '%H:%M').time()
    rec.nap_duration = nap_duration
    rec.efficiency = efficiency
    session.commit()
    session.close()

def get_sleep_records(days=30):
    from models import SleepRecord
    session = Session()
    since = date.today() - timedelta(days=days)
    res = session.query(SleepRecord).filter(SleepRecord.date >= since).all()
    session.close()
    return res

STATE_ADVICE_MAP = {'犯困':'建议起身活动5分钟，用冷水洗脸，或听一首快节奏歌曲提神。','走神':'尝试番茄工作法：专注25分钟，休息5分钟。当前可先休息2分钟深呼吸。','疲惫':'建议闭眼休息10分钟，或者吃点零食补充能量。','高效':'保持当前状态，建议继续攻克重难点。','正常':'状态平稳，按计划进行。'}

def add_state(state, advice=''):
    from models import StateEvent
    if not advice:
        advice = STATE_ADVICE_MAP.get(state, '保持节奏，适当休息。')
    session = Session()
    event = StateEvent(state=state, advice=advice)
    session.add(event)
    session.commit()
    session.close()
    return advice

def get_states(days=30):
    from models import StateEvent
    session = Session()
    since = datetime.now() - timedelta(days=days)
    res = session.query(StateEvent).filter(StateEvent.datetime >= since).all()
    session.close()
    return res

def get_state_stats(days=30):
    states = get_states(days)
    if not states:
        return {}
    from collections import Counter
    counter = Counter([s.state for s in states])
    hour_counts = {}
    for s in states:
        h = s.datetime.hour
        hour_counts[h] = hour_counts.get(h, 0) + 1
    return {'counter': dict(counter), 'hourly': hour_counts}

def add_events_bulk(events_list):
    from models import StudyEvent, KnowledgeNode
    from config import REVIEW_INTERVALS
    session = Session()
    for ev in events_list:
        dt = datetime.strptime(ev['date'], '%Y-%m-%d').date()
        event = StudyEvent(date=dt, subject=ev['subject'], topic=ev['topic'], activity=ev['activity'], duration=int(ev['duration']), note=ev.get('note', ''))
        session.add(event)
        node = session.query(KnowledgeNode).filter_by(subject=ev['subject'], topic=ev['topic']).first()
        if not node:
            node = KnowledgeNode(subject=ev['subject'], topic=ev['topic'], last_review=dt)
            session.add(node)
        node.last_review = dt
        for interval in REVIEW_INTERVALS:
            candidate = dt + timedelta(days=interval)
            if candidate > dt:
                node.next_review = candidate
                break
        else:
            node.next_review = dt + timedelta(days=REVIEW_INTERVALS[-1])
    session.commit()
    session.close()
    from trainer import train_model

def clear_all_data():
    from models import StudyEvent, KnowledgeNode, UserGoal, DailyPlan, SleepRecord, StateEvent, Task, KnowledgePoint
    session = Session()
    session.query(StudyEvent).delete()
    session.query(KnowledgeNode).delete()
    session.query(UserGoal).delete()
    session.query(DailyPlan).delete()
    session.query(SleepRecord).delete()
    session.query(StateEvent).delete()
    session.query(Task).delete()
    session.query(KnowledgePoint).delete()
    session.commit()
    session.close()

def delete_event(event_id):
    from models import StudyEvent
    session = Session()
    event = session.query(StudyEvent).filter_by(id=event_id).first()
    if event:
        session.delete(event)
    session.commit()
    session.close()
    from trainer import train_model

def parse_document(filepath):
    events = []
    text = ""
    ext = os.path.splitext(filepath)[1].lower()
    try:
        if ext == '.docx':
            from docx import Document
            doc = Document(filepath)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif ext == '.pdf':
            import PyPDF2
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            return []
    except Exception as e:
        raise Exception(f"解析文件失败：{e}")
    lines = text.splitlines()
    for line in lines:
        line = line.strip()
        if not line:
            continue
        date_match = re.search(r'(\d{4}-\d{2}-\d{2})', line)
        if not date_match:
            continue
        date_str = date_match.group(1)
        subject = None
        for sub in ['数学','英语','专业课','政治']:
            if sub in line:
                subject = sub
                break
        if not subject:
            continue
        activity = None
        for act in ['新学','复习','做题','看视频']:
            if act in line:
                activity = act
                break
        if not activity:
            activity = '新学'
        dur_match = re.search(r'(\d+)\s*分钟', line)
        if not dur_match:
            continue
        duration = int(dur_match.group(1))
        topic = line
        for token in [date_str, subject, activity, f"{duration}分钟"]:
            topic = topic.replace(token, '')
        topic = topic.strip(' -：,，').strip()
        if not topic:
            topic = "未命名主题"
        events.append({'date': date_str, 'subject': subject, 'topic': topic, 'activity': activity, 'duration': str(duration), 'note': ''})
    return events

def add_task(subject, phase, description, plan_date_str):
    from models import Task
    session = Session()
    task = Task(subject=subject, phase=phase, description=description, plan_date=datetime.strptime(plan_date_str, '%Y-%m-%d').date())
    session.add(task)
    session.commit()
    session.close()

def get_tasks(subject=None, phase=None, status=None):
    from models import Task
    session = Session()
    q = session.query(Task)
    if subject:
        q = q.filter_by(subject)
    if phase:
        q = q.filter_by(phase)
    if status:
        q = q.filter_by(status)
    res = q.all()
    session.close()
    return res

def update_task_status(task_id, new_status, actual_date=None):
    from models import Task
    session = Session()
    task = session.query(Task).filter_by(id=task_id).first()
    if task:
        task.status = new_status
        if new_status == '已完成' and actual_date:
            task.actual_date = datetime.strptime(actual_date, '%Y-%m-%d').date()
        elif new_status == '已完成' and not actual_date:
            task.actual_date = date.today()
        session.commit()
    session.close()

def delete_task(task_id):
    from models import Task
    session = Session()
    task = session.query(Task).filter_by(id=task_id).first()
    if task:
        session.delete(task)
    session.commit()
    session.close()

def add_knowledge_point(subject, title, status='未学', priority=0, importance=0.0, source=''):
    from models import KnowledgePoint
    session = Session()
    kp = KnowledgePoint(subject=subject, title=title, status=status, priority=priority, importance=importance, source=source)
    session.add(kp)
    session.commit()
    session.close()

def add_knowledge_points_bulk(kps_list):
    from models import KnowledgePoint
    session = Session()
    for item in kps_list:
        kp = KnowledgePoint(subject=item['subject'], title=item['title'], priority=item.get('priority', 0), importance=item.get('importance', 0.0))
        session.add(kp)
    session.commit()
    session.close()

def get_knowledge_points(subject=None, status=None):
    from models import KnowledgePoint
    session = Session()
    q = session.query(KnowledgePoint)
    if subject:
        q = q.filter_by(subject)
    if status:
        q = q.filter_by(status)
    res = q.all()
    session.close()
    return res

def update_knowledge_point_status(title, new_status):
    from models import KnowledgePoint
    session = Session()
    kp = session.query(KnowledgePoint).filter_by(title=title).first()
    if kp:
        kp.status = new_status
    session.commit()
    session.close()

def get_all_study_records():
    from models import StudyEvent
    session = Session()
    records = session.query(StudyEvent).order_by(StudyEvent.date.desc(), StudyEvent.id.desc()).all()
    res = []
    for r in records:
        res.append({
            "id": r.id,
            "date": r.date,
            "subject": r.subject,
            "topic": r.topic,
            "activity": r.activity,
            "duration": r.duration,
            "note": r.note
        })
    session.close()
    return res

def get_all_sleep_records():
    from models import SleepRecord
    session = Session()
    records = session.query(SleepRecord).order_by(SleepRecord.date.desc(), SleepRecord.id.desc()).all()
    res = []
    for r in records:
        res.append({
            "id": r.id,
            "date": r.date,
            "sleep_time": r.bed_time,
            "wake_time": r.wake_time,
            "nap_duration": r.nap_duration
        })
    session.close()
    return res

def get_all_daily_tasks():
    from models import DailyTask
    session = Session()
    tasks = session.query(DailyTask).order_by(DailyTask.date.desc(), DailyTask.id.desc()).all()
    res = []
    for t in tasks:
        res.append({
            "id": t.id,
            "date": t.date,
            "content": t.content,
            "completed": t.completed
        })
    session.close()
    return res

def get_all_stage_tasks():
    from models import StageTask
    session = Session()
    tasks = session.query(StageTask).order_by(StageTask.id.desc()).all()
    res = []
    for t in tasks:
        res.append({
            "id": t.id,
            "title": t.title,
            "deadline": t.deadline,
            "progress": t.progress,
            "completed": t.completed
        })
    session.close()
    return res

def get_all_status_reports():
    from models import StatusReport
    session = Session()
    reports = session.query(StatusReport).order_by(StatusReport.id.desc()).all()
    res = []
    for r in reports:
        res.append({
            "id": r.id,
            "date": r.date,
            "content": r.content,
            "feedback": r.feedback
        })
    session.close()
    return res

def delete_all_knowledge_points():
    from models import KnowledgePoint
    session = Session()
    session.query(KnowledgePoint).delete()
    session.commit()
    session.close()

def parse_knowledge_document(filepath):
    text = ""
    ext = os.path.splitext(filepath)[1].lower()
    try:
        if ext == '.txt':
            with open(filepath, 'r', encoding='utf-8') as f:
                text = f.read()
        elif ext == '.csv':
            with open(filepath, 'r', encoding='utf-8') as f:
                reader = csv.reader(f)
                headers = next(reader, None)
                kps = []
                for row in reader:
                    if len(row) >= 2:
                        kps.append({'subject': row[0].strip(), 'title': row[1].strip(), 'priority': int(row[2]) if len(row) > 2 and row[2].isdigit() else 0, 'importance': float(row[3]) if len(row) > 3 and row[3].replace('.', '', 1).isdigit() else 0.0})
                return kps
        elif ext in ['.docx', '.pdf', '.pptx']:
            if ext == '.docx':
                from docx import Document
                doc = Document(filepath)
                for para in doc.paragraphs:
                    text += para.text + "\n"
            elif ext == '.pdf':
                import PyPDF2
                with open(filepath, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
            elif ext == '.pptx':
                from pptx import Presentation
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
        else:
            return []
    except Exception as e:
        raise Exception(f"解析文件失败：{e}")
    kps = []
    lines = text.splitlines()
    for line in lines:
        line = line.strip()
        if not line:
            continue
        if '：' in line:
            parts = line.split('：', 1)
        elif ':' in line:
            parts = line.split(':', 1)
        else:
            parts = re.split(r'\s+', line, 1)
        if len(parts) == 2:
            subject = parts[0].strip()
            title = parts[1].strip()
            if subject in ['数学','英语','专业课','政治'] and title:
                kps.append({'subject': subject, 'title': title, 'priority': 0, 'importance': 0.0})
    return kps
